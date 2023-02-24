# -*- coding: utf-8 -*-
'''
  スクリプト名 : py_powerpoint_20230224.py
  処理内容：powerpointファイルを作成
  補足： pip install python-pptx
        公式ライブラリは https://openpyxl.readthedocs.io/en/stable/
'''

# 必要なモジュールたちを読み込み
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt, Cm # Python-pptxでのデフォルト単位は、emu(English Metric Units)という単位→ cmに変換する
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import time

from matplotlib import pyplot as plt
import pandas as pd
import datetime
import os
import sys


if __name__ == "__main__":
   # カレントディレクトリを起動パスのディレクトリに変更
   os.chdir(os.path.dirname(os.path.abspath(sys.argv[0])))
   
   # 現在時刻
   now_value = datetime.datetime.now()
   
  
   # Presentationオブジェクトの宣言
   presentation = Presentation()

  
   # 表紙スライドのレイアウトを取得
   slide_layout1 = presentation.slide_layouts[0] # レイアウト番号
   
   # 表紙スライドを作成
   slide = presentation.slides.add_slide(slide_layout1) # add_slide() でスライドを1枚追加
   
   # スライドのタイトル情報を取得して更新
   title = slide.shapes.title
   title.text = "PythonでPowerPoint"
   
   # サブタイトル部分の情報を取得して更新
   subtitle = slide.placeholders[1]
   subtitle.text = "作成者: tommyda"
   
   # 次ページ
   slide_layout1 = presentation.slide_layouts[0] # レイアウト番号
   slide = presentation.slides.add_slide(slide_layout1)
   
   title = slide.shapes.title
   # タイトルを入力
   title.text = "PythonでPowerPointを作ります" 

   
   # 本日の日付を作成
   today = now_value.strftime('%Y年%m月%d日')
   
   subtitle = slide.placeholders[1]
   subtitle.text = "作成者: tommyda\n" + today # \n は改行を意味する   

   # 次ページ
   slide_layout2 = presentation.slide_layouts[1] # レイアウト番号
   slide = presentation.slides.add_slide(slide_layout2)
   
   title = slide.shapes.title
   content = slide.placeholders[1] # スライド内の、2つ目の枠
   
   title.text = "目次"
   content.text = "テスト１\nテスト２\nテスト３"
   
   # 次ページ
   slide_layout3 = presentation.slide_layouts[1]
   slide = presentation.slides.add_slide(slide_layout3)
   
   title = slide.shapes.title
   content = slide.placeholders[1]
   
   # 画像挿入
   image_file_name = "apple_icon.jpg"
   
   title.text = "リンゴのサンプル"
   # 画像ファイル名/左上のX座標/左YのX座標/画像の幅/画像の高さ　を指定する
   pic = slide.shapes.add_picture(image_file_name, Cm(10), Cm(10), Cm(7), Cm(7)) 
   
   
   slide_layout2 = presentation.slide_layouts[1]
   slide = presentation.slides.add_slide(slide_layout2)
   
   title = slide.shapes.title
   content = slide.placeholders[1]
   
   title.text = "表を入れる"
   content.text = "テスト"
   
   rows = 3 # 行数
   cols = 3 # 列数
   
   # 表を追加。行数/列数/左上のX座標/左YのX座標/表の幅/表の高さ を指定する
   table_shape = slide.shapes.add_table(rows, cols, Cm(3), Cm(10), Cm(20), Cm(8))
   table = table_shape.table
   
   table.cell(0, 0).text = "A" # 1行1列目
   table.cell(0, 1).text = "B" # 1行2列目
   table.cell(0, 2).text = "C" # 1行3列目
   table.cell(1, 0).text = "D" # 2行1列目
   table.cell(1, 1).text = "E" # 2行2列目
   table.cell(1, 2).text = "F" # 2行3列目
   table.cell(2, 0).text = "G" # 3行1列目
   table.cell(2, 1).text = "H" # 3行2列目
   table.cell(2, 2).text = "I" # 3行3列目
   
   
   # 新しいスライドを追加する
   slide_layout = presentation.slide_layouts[5]
   slide = presentation.slides.add_slide(slide_layout)
   
   title = slide.shapes.title
   title.text = "グラフを入れる"
   
   # スライドにグラフを挿入する
   chart_data = CategoryChartData()
   chart_data.categories = ['A', 'B', 'C']
   chart_data.add_series('売上高', (34, 25, 27))
   x, y, cx, cy = presentation.slide_width // 4, presentation.slide_height // 4, presentation.slide_width // 2, presentation.slide_height // 2
   chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
   
   # ファイル名
   nowtime = now_value.strftime("%Y%m%d%H%M%S")
   file_name = 'python_powerpoint_{}.pptx'.format(nowtime)
   
  
   # 保存
   presentation.save(file_name)
   
   # 開始時刻からの経過時間を取得
   dtend = datetime.datetime.now()
   delta = dtend - now_value
   print("処理終了 ", dtend, " 経過時間 :  ", delta)   